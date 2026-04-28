from __future__ import annotations

import base64
import mimetypes
import os
import re
import shutil
import subprocess
import uuid
from pathlib import Path
from typing import Any, TypedDict

from openai import OpenAI
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from langchain_core.messages import HumanMessage, SystemMessage
from langchain_openai import ChatOpenAI
from langchain_tavily import TavilySearch
from langgraph.graph import END, START, StateGraph

from config import DEFAULT_STYLE, DEFAULT_TONE, DEFAULT_VOICE, require_env


LLM_MODEL = "gpt-4o-mini"
TTS_MODEL = "gpt-4o-mini-tts"
REQUIRED_SYSTEM_TOOLS = ("ffmpeg", "ffprobe", "soffice", "xvfb-run", "pdftoppm")


class State(TypedDict, total=False):
    pptx_path: str
    work_dir: str
    prompt: dict[str, Any]
    slides: list[dict[str, Any]]
    n_slides: int
    slide_index: int
    audios: list[str]
    videos: list[str]
    page_contents: list[str]
    scripts: list[str]
    slide_images: list[str]
    final_video: str
    summary: str
    summary_img: str


def ensure_runtime_requirements() -> None:
    missing_tools = [tool for tool in REQUIRED_SYSTEM_TOOLS if shutil.which(tool) is None]
    if missing_tools:
        raise RuntimeError(
            "필수 시스템 도구가 없습니다: "
            + ", ".join(missing_tools)
            + ". README의 런타임 설치 단계를 먼저 수행하세요."
        )

    require_env("OPENAI_API_KEY")
    require_env("TAVILY_API_KEY")


def clean_text(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip()


def ffprobe_duration(path: str) -> float:
    out = subprocess.check_output(
        [
            "ffprobe",
            "-v",
            "error",
            "-show_entries",
            "format=duration",
            "-of",
            "default=noprint_wrappers=1:nokey=1",
            path,
        ]
    ).decode().strip()
    return float(out)


def img_to_data_url(path: str) -> str:
    mime = mimetypes.guess_type(path)[0] or "image/png"
    with open(path, "rb") as file:
        encoded = base64.b64encode(file.read()).decode("utf-8")
    return f"data:{mime};base64,{encoded}"


def render_mp4(image_path: str, audio_path: str, out_mp4: str, width: int = 1920, height: int = 1080) -> None:
    duration = ffprobe_duration(audio_path)
    vf = (
        f"scale={width}:{height}:force_original_aspect_ratio=decrease,"
        f"pad={width}:{height}:(ow-iw)/2:(oh-ih)/2:color=black"
    )
    cmd = [
        "ffmpeg",
        "-y",
        "-loop",
        "1",
        "-i",
        image_path,
        "-i",
        audio_path,
        "-t",
        str(duration),
        "-vf",
        vf,
        "-c:v",
        "libx264",
        "-threads",
        "1",
        "-preset",
        "veryfast",
        "-crf",
        "20",
        "-c:a",
        "aac",
        "-b:a",
        "192k",
        "-pix_fmt",
        "yuv420p",
        "-movflags",
        "+faststart",
        out_mp4,
    ]
    subprocess.check_call(cmd)


def make_summary_video(summary_text: str, audio_path: str, out_path: str) -> None:
    # Try to find a common Korean font on Windows and Linux
    font_paths = [
        "/usr/share/fonts/truetype/nanum/NanumBarunGothic.ttf", # Linux Docker
        "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",      # Linux Docker
        "/usr/share/fonts/truetype/nanum/NanumGothicCoding.ttf",# Linux Docker
        "C:/Windows/Fonts/malgun.ttf",  # Malgun Gothic
        "C:/Windows/Fonts/batang.ttc",  # Batang
        "C:/Windows/Fonts/gulim.ttc",   # Gulim
    ]
    font_path = ""
    for p in font_paths:
        if os.path.exists(p):
            font_path = p
            break

    # If no specific font found, let ffmpeg try to find one or use default
    font_clause = f"fontfile='{font_path}':" if font_path else ""

    work_dir = os.path.dirname(out_path)
    text_file_path = os.path.join(work_dir, "summary_text.txt")
    with open(text_file_path, "w", encoding="utf-8") as f:
        f.write(summary_text)

    # 1. Create a 1920x1080 black image first (extremely light for memory)
    black_img_path = os.path.join(work_dir, "black_bg.png")
    subprocess.run([
        "ffmpeg", "-y", "-f", "lavfi", "-i", "color=c=black:s=1920x1080",
        "-frames:v", "1", black_img_path
    ], check=True)

    duration = ffprobe_duration(audio_path)
    escaped_text_file_path = text_file_path.replace("\\", "/").replace(":", "\\:")

    # 2. Render summary video using the black image (more memory efficient than lavfi)
    cmd = [
        "ffmpeg",
        "-y",
        "-loop", "1",
        "-i", black_img_path,
        "-i", audio_path,
        "-t", str(duration),
        "-vf",
        (
            f"drawtext={font_clause}"
            f"textfile='{escaped_text_file_path}':"
            "fontcolor=white:fontsize=60:x=(w-text_w)/2:y=(h-text_h)/2:"
            "line_spacing=30"
        ),
        "-c:v", "libx264",
        "-threads", "1",
        "-preset", "veryfast",
        "-crf", "23",
        "-c:a", "aac",
        "-b:a", "128k",
        "-pix_fmt", "yuv420p",
        out_path,
    ]

    subprocess.run(cmd, check=True)


def export_slide_as_png(state: dict[str, Any], dpi: int = 150) -> dict[str, Any]:
    work_dir = Path(state["work_dir"]).expanduser().resolve()
    work_dir.mkdir(parents=True, exist_ok=True)

    pptx_path = Path(state["pptx_path"]).expanduser().resolve()
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX 없음: {pptx_path}")

    page_no = int(state.get("slide_index", 0)) + 1
    out_prefix = work_dir / "slide_img"
    lo_profile = f"file:///tmp/lo_profile_{uuid.uuid4().hex}"

    env = os.environ.copy()
    env.update(
        {
            "LANG": "ko_KR.UTF-8",
            "LC_ALL": "ko_KR.UTF-8",
            "FONTCONFIG_PATH": str(Path("~/.config/fontconfig").expanduser()),
            "FONTCONFIG_FILE": str(Path("~/.config/fontconfig/fonts.conf").expanduser()),
            "HOME": str(Path("~").expanduser()),
            "XDG_CACHE_HOME": str(Path("~/.cache").expanduser()),
            "SAL_USE_VCLPLUGIN": "gen",
        }
    )

    def run_lo_convert(convert_to: str) -> subprocess.CompletedProcess[str]:
        cmd = [
            "xvfb-run",
            "-a",
            "soffice",
            "--headless",
            "--nologo",
            "--nofirststartwizard",
            "--norestore",
            f"-env:UserInstallation={lo_profile}",
            "--convert-to",
            convert_to,
            "--outdir",
            str(work_dir),
            str(pptx_path),
        ]
        return subprocess.run(cmd, capture_output=True, text=True, env=env, check=False)

    before_png = set(work_dir.glob("*.png"))
    run_lo_convert("png:impress_png_Export")
    created_png = [path for path in work_dir.glob("*.png") if path not in before_png]

    exact = [path for path in created_png if path.stem.endswith(f"-{page_no}")]
    candidate = None
    if exact:
        candidate = max(exact, key=lambda path: path.stat().st_mtime)
    elif created_png:
        candidate = max(created_png, key=lambda path: path.stat().st_mtime)

    if candidate and candidate.exists():
        state["slide_image"] = str(candidate)
        return state

    before_pdf = set(work_dir.glob("*.pdf"))
    res_pdf = run_lo_convert("pdf:impress_pdf_Export")
    target_pdf = work_dir / f"{pptx_path.stem}.pdf"
    created_pdf = [path for path in work_dir.glob("*.pdf") if path not in before_pdf]

    if target_pdf.exists():
        pdf_path = target_pdf
    elif created_pdf:
        pdf_path = max(created_pdf, key=lambda path: path.stat().st_mtime)
    else:
        raise RuntimeError(
            "PPTX -> PDF 변환 실패\n"
            f"stdout: {res_pdf.stdout}\n"
            f"stderr: {res_pdf.stderr}"
        )

    ppm_cmd = [
        "pdftoppm",
        "-f",
        str(page_no),
        "-l",
        str(page_no),
        "-png",
        "-r",
        str(dpi),
        "-scale-to",
        "1920",
        str(pdf_path),
        str(out_prefix),
    ]
    res2 = subprocess.run(ppm_cmd, capture_output=True, text=True, env=env, check=False)

    png_path = Path(f"{out_prefix}-{page_no}.png")
    if not png_path.exists():
        raise RuntimeError(
            "PDF -> PNG 변환 실패\n"
            f"stdout: {res2.stdout}\n"
            f"stderr: {res2.stderr}"
        )

    state["slide_image"] = str(png_path)
    return state


def concat_videos_ffmpeg(video_paths: list[str], out_path: str, reencode: bool = False) -> None:
    list_path = out_path + ".txt"
    with open(list_path, "w", encoding="utf-8") as file:
        for video_path in video_paths:
            file.write(f"file '{os.path.abspath(video_path)}'\n")

    if reencode:
        # Robust concatenation: Ensure all streams are scaled/padded to 1920x1080 yuv420p before encoding
        # This prevents resolution mismatches between slides and summary.
        cmd = [
            "ffmpeg",
            "-y",
            "-safe", "0",
            "-f", "concat",
            "-i", list_path,
            "-vf", "scale=1920:1080:force_original_aspect_ratio=decrease,pad=1920:1080:(ow-iw)/2:(oh-ih)/2,format=yuv420p",
            "-c:v", "libx264",
            "-threads", "1",
            "-preset", "veryfast",
            "-c:a", "aac",
            "-b:a", "192k",
            out_path,
        ]
    else:
        cmd = ["ffmpeg", "-y", "-safe", "0", "-f", "concat", "-i", list_path, "-c", "copy", out_path]

    subprocess.check_call(cmd)


def build_llm() -> ChatOpenAI:
    return ChatOpenAI(model=LLM_MODEL, temperature=0.3)


def node_parse_all(state: State) -> State:
    presentation = Presentation(state["pptx_path"])
    work_dir = state["work_dir"]
    media_dir = os.path.join(work_dir, "media")
    os.makedirs(media_dir, exist_ok=True)

    slides_out: list[dict[str, Any]] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        tables: list[list[list[str]]] = []
        images: list[str] = []
        title = ""

        if slide.shapes.title:
            title = clean_text(slide.shapes.title.text)

        for shape_index, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs)
                text = clean_text(text)
                if text:
                    texts.append(text)
                    if not title and shape_index == 0:
                        title = text[:30]

            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = [[clean_text(cell.text) for cell in row.cells] for row in shape.table.rows]
                tables.append(table)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                ext = shape.image.ext
                image_path = os.path.join(media_dir, f"slide{idx}_img_{shape_index}.{ext}")
                with open(image_path, "wb") as file:
                    file.write(shape.image.blob)
                images.append(image_path)

        snap_state = export_slide_as_png(
            {"pptx_path": state["pptx_path"], "work_dir": work_dir, "slide_index": idx - 1}
        )
        slides_out.append(
            {
                "index": idx,
                "title": title,
                "texts": texts,
                "tables": tables,
                "images": images,
                "snap": snap_state["slide_image"],
            }
        )

    state["slides"] = slides_out
    state["n_slides"] = len(slides_out)
    state["slide_index"] = 0
    state["page_contents"] = []
    state["scripts"] = []
    state["audios"] = []
    state["videos"] = []
    state["slide_images"] = []
    state["summary"] = ""
    return state


def node_tool_search(state: State) -> State:
    """실제 검색을 수행하고 요약된 문자열을 각 슬라이드에 반환하는 노드"""
    slides = state.get("slides", [])
    search = TavilySearch(max_results=3)

    for slide in slides:
        title = slide.get("title", "")
        if title:
            try:
                slide["search_context"] = str(search.run(title))
            except Exception as exc:
                slide["search_context"] = f"검색 중 오류 발생: {exc}"
        else:
            slide["search_context"] = "검색할 제목이 없어 검색을 수행하지 않았습니다."

    return state


def node_gen_page(state: State) -> State:
    """슬라이드 정보를 종합하여 핵심 요약을 작성하는 노드 (일괄 처리)"""
    slides = state.get("slides", [])
    prompt_config = state["prompt"]
    state["page_contents"] = []

    for slide in slides:
        texts = slide.get("texts", [])
        tables = slide.get("tables", [])
        images = slide.get("images", [])
        snap = slide.get("snap")
        search_info = slide.get("search_context", "관련 외부 정보 없음")

        table_snip = ""
        if tables:
            try:
                table_snip = "\n".join(" | ".join(map(str, row)) for row in tables[0][:6])
            except Exception:
                table_snip = str(tables[0][:6])

        combined_images: list[str] = []
        if snap:
            combined_images.append(snap)
        for image in images:
            if image not in combined_images:
                combined_images.append(image)

        system_msg = SystemMessage(
            content=f"""
당신은 전문적인 AI 강사입니다.
슬라이드의 텍스트, 표, 이미지를 종합하여 슬라이드 핵심 요약을 작성하세요.

아래 작성 규칙을 준수하여 작성하세요
[작성 규칙]
1. 단락 구성: 반드시 4문장에서 6문장 사이의 하나의 단락으로 작성하세요.
2. 불릿 금지: 나열식 기호( -, *, • )나 숫자를 절대 사용하지 마세요. 모든 내용은 서술형 문장으로 이어져야 합니다.
3. 이미지 통합: 첨부된 이미지에서 보이는 시각적 정보나 도표의 특징을 요약 내용에 자연스럽게 포함하세요.
4. 객관성: 추측이나 과장된 표현을 배제하고, 제공된 데이터에만 기반하여 사실적으로 작성하세요.
5. 요약 스타일: {prompt_config.get("style", DEFAULT_STYLE)}
""".strip()
        )

        user_text = (
            f"[텍스트]\n{texts if texts else '(없음)'}\n\n"
            f"[표 요약]\n{table_snip if table_snip else '(없음)'}\n\n"
            f"[외부 참고 지식]\n{search_info}"
        )
        user_content: list[dict[str, Any]] = [{"type": "text", "text": user_text}]

        for image_path in combined_images[:3]:
            try:
                user_content.append({"type": "image_url", "image_url": {"url": img_to_data_url(image_path)}})
            except Exception:
                pass

        response = build_llm().invoke([system_msg, HumanMessage(content=user_content)])
        state["page_contents"].append(str(response.content).strip())

    return state


def node_gen_script(state: State) -> State:
    idx = state.get("slide_index", 0)
    slides = state.get("slides", [])
    n_slides = state.get("n_slides", len(slides))

    if idx >= len(slides):
        return state

    page_content = state.get("page_contents", [])[idx]
    prompt_config = state.get("prompt", {})
    work_dir = state.get("work_dir", ".")

    is_first = (idx == 0)
    is_last = (idx == n_slides - 1)

    # 이전 슬라이드 대본 가져오기
    all_scripts = state.get("scripts", [])
    if is_first:
        previous_script = "이전 슬라이드 없음"
    else:
        previous_script = all_scripts[idx - 1] if idx - 1 < len(all_scripts) else ""
        previous_script = previous_script[-200:] if len(previous_script) > 200 else previous_script

    # 마지막 슬라이드에서는 다음 슬라이드 요약을 넘기지 않음
    if is_last:
        next_context = "현재 슬라이드는 마지막 슬라이드입니다. 다음 슬라이드가 없습니다."
    elif idx + 1 < len(state.get("page_contents", [])):
        raw_next = state["page_contents"][idx + 1]
        next_context = raw_next[:200] if len(raw_next) > 200 else raw_next
    else:
        next_context = "다음 슬라이드 요약 없음"

    system_msg = SystemMessage(
        content=f"""
당신은 대학 강의 발표 대본 작성 보조 에이전트입니다.
슬라이드 요약을 바탕으로 60~90초 분량의 발표 대본을 작성하세요.

[현재 슬라이드 정보]
- 현재 슬라이드 번호: {idx + 1}
- 전체 슬라이드 수: {n_slides}
- 첫 번째 슬라이드 여부: {is_first}
- 마지막 슬라이드 여부: {is_last}

[작성 지침]
현재 슬라이드 1장에 해당하는 발표 스크립트만 작성하세요.
전체 발표 흐름은 고려하되, 현재 슬라이드의 핵심 내용이 중심이 되어야 합니다.
이전 스크립트는 자연스러운 연결을 위한 참고용이며, 이전 문장을 반복하지 마세요.
슬라이드에 없는 내용을 임의로 추가하지 마세요.

첫 번째 슬라이드라면 반드시 간단한 인사로 시작하세요.
예: "안녕하세요. 오늘은 ~에 대해 이야기해보겠습니다."

두 번째 이후 슬라이드라면 인사말을 사용하지 마세요.
두 번째 이후 슬라이드는 이전 내용에서 현재 내용으로 자연스럽게 연결하세요.

마지막 슬라이드라면 발표 전체를 자연스럽게 마무리하세요.
마지막 슬라이드에서는 절대 다음 슬라이드, 다음 내용, 이후 내용을 언급하지 마세요.

마지막 슬라이드가 아니라면 "감사합니다", "이상입니다", "마치겠습니다" 같은 종료 표현을 사용하지 마세요.

[금지 표현]
- 다음 슬라이드에서는
- 다음 슬라이드에서
- 다음으로는
- 이후에는
- 뒤에서 설명하겠습니다
- 이어서
- 계속해서

[규칙]
- 불릿 기호(-, •, *) 사용 금지
- 번호 목록 사용 금지
- 코드블록 사용 금지
- 자연스러운 구어체로 작성
- '~습니다', '~합니다' 체 사용
- 숫자와 용어는 원문 기준으로 정확하게 유지
- 이전 슬라이드와 같은 문장 구조를 반복하지 않기
- 말투/톤: {prompt_config.get('tone', DEFAULT_TONE)}
- 스타일: {prompt_config.get('style', DEFAULT_STYLE)}
""".strip()
    )

    user_msg = HumanMessage(
        content=f"""
[이전 슬라이드 대본 맥락]
{previous_script}

[현재 슬라이드 요약]
{page_content}

[다음 슬라이드 참고 정보]
{next_context}

주의:
다음 슬라이드 참고 정보는 전체 흐름 파악용입니다.
대본 안에서 다음 슬라이드를 예고하거나 언급하지 마세요.
""".strip()
    )

    response = build_llm().invoke([system_msg, user_msg])
    script = str(response.content).strip()

    bad_phrases = [
        "다음 슬라이드에서는",
        "다음 슬라이드에서",
        "다음으로는",
        "이후에는",
        "뒤에서 설명하겠습니다",
        "계속해서",
    ]
    for phrase in bad_phrases:
        script = script.replace(phrase, "")

    script_path = os.path.join(work_dir, f"script_{idx + 1}.txt")
    with open(script_path, "w", encoding="utf-8") as file:
        file.write(script)

    state["scripts"].append(script)
    return state


def node_tts(state: State) -> State:
    client = OpenAI()
    idx = state.get("slide_index", 0)
    script = state["scripts"][idx]
    response = client.audio.speech.create(
        model=TTS_MODEL,
        voice=state["prompt"].get("voice", DEFAULT_VOICE),
        input=script,
    )
    mp3_path = os.path.join(state["work_dir"], f"narration_{idx + 1}.mp3")
    with open(mp3_path, "wb") as file:
        file.write(response.content)
    state["audios"].append(mp3_path)
    return state


def node_make_video(state: State) -> State:
    idx = state.get("slide_index", 0)
    slide = state["slides"][idx]
    slide_image = slide.get("snap")
    audio_path = state["audios"][idx]

    if not slide_image:
        raise ValueError(f"Slide {idx + 1}의 이미지 경로가 없습니다.")
    if not audio_path:
        raise ValueError(f"Slide {idx + 1}의 오디오 경로가 없습니다.")

    output_path = os.path.join(state["work_dir"], f"slide_{idx + 1}_video.mp4")
    render_mp4(slide_image, audio_path, output_path)
    state["videos"].append(output_path)
    state["slide_index"] = idx + 1
    return state


def node_summary(state: State) -> State:
    """전체 내용을 요약하고 '3줄 핵심 요약' 영상을 제작하는 노드"""
    page_contents = state.get("page_contents", [])
    if not page_contents:
        return state

    all_contents = "\n".join(page_contents)

    summary_prompt = f"""
다음은 전체 강의 내용입니다.
마지막에 시청자에게 전달할 '핵심 요약 3줄'을 강의 톤에 맞춰 작성해주세요.

[강의 내용]
{all_contents}

[규칙]
- 반드시 "마지막으로 오늘 배운 내용을 요약해 보겠습니다"라는 멘트로 시작할 것.
- 가장 중요한 포인트 3가지만 간결하게 설명할 것.
- 반드시 마지막에는 "이상으로 발표를 마칩니다. 들어주셔서 감사합니다."
""".strip()

    summary_img_prompt = f"""
다음 강의 내용을 핵심만 3줄로 요약해줘.
각 줄은 반드시 20자 내외로 작성하고, 각 포인트 앞에 숫자를 붙여줘.

[강의 내용]
{all_contents}

[출력 형식 예시]
1. 인공지능의 정의와 역사 이해
2. 딥러닝 모델의 학습 원리 파악
3. 실무 적용을 위한 파이프라인 구축
""".strip()

    llm = build_llm()
    summary_script = llm.invoke(summary_prompt).content
    summary_img = llm.invoke(summary_img_prompt).content.strip()

    # 요약용 음성(TTS) 생성
    summary_mp3_path = os.path.join(state["work_dir"], "summary_audio.mp3")
    client = OpenAI()
    resp = client.audio.speech.create(
        model=TTS_MODEL,
        voice=state["prompt"].get("voice", DEFAULT_VOICE),
        input=summary_script,
    )
    with open(summary_mp3_path, "wb") as f:
        f.write(resp.content)

    # 요약 영상 제작
    summary_video_path = os.path.join(state["work_dir"], "summary_video.mp4")
    make_summary_video(summary_img, summary_mp3_path, summary_video_path)

    state["scripts"].append(summary_script)
    state["videos"].append(summary_video_path)
    state["summary"] = summary_script
    state["summary_img"] = summary_img

    return state


def node_concat(state: State) -> State:
    output_path = os.path.join(state["work_dir"], "final_lecture_video.mp4")
    concat_videos_ffmpeg(state["videos"], output_path, reencode=True)
    state["final_video"] = output_path
    return state


def decide_next_step(state: State) -> str:
    return "CONTINUE" if state.get("slide_index", 0) < state.get("n_slides", 0) else "END"


def build_graph():
    builder = StateGraph(State)
    builder.add_node("parse_all", node_parse_all)
    builder.add_node("tool_search", node_tool_search)
    builder.add_node("gen_page", node_gen_page)
    builder.add_node("gen_script_ctx", node_gen_script)
    builder.add_node("tts", node_tts)
    builder.add_node("make_video", node_make_video)
    builder.add_node("summary", node_summary)
    builder.add_node("concat", node_concat)

    builder.add_edge(START, "parse_all")
    builder.add_edge("parse_all", "tool_search")
    builder.add_edge("tool_search", "gen_page")
    builder.add_edge("gen_page", "gen_script_ctx")
    builder.add_edge("gen_script_ctx", "tts")
    builder.add_edge("tts", "make_video")
    builder.add_conditional_edges("make_video", decide_next_step, {"CONTINUE": "gen_script_ctx", "END": "summary"})
    builder.add_edge("summary", "concat")
    builder.add_edge("concat", END)
    return builder.compile()


GRAPH = build_graph()


def build_initial_state(pptx_path: str, work_dir: str, tone: str | None, voice: str | None, style: str | None) -> State:
    return {
        "pptx_path": pptx_path,
        "work_dir": work_dir,
        "prompt": {
            "voice": voice or DEFAULT_VOICE,
            "tone": tone or DEFAULT_TONE,
            "style": style or DEFAULT_STYLE,
        },
    }


def run_pipeline(pptx_path: str, work_dir: str, tone: str | None = None, voice: str | None = None, style: str | None = None) -> State:
    ensure_runtime_requirements()
    state = build_initial_state(pptx_path=pptx_path, work_dir=work_dir, tone=tone, voice=voice, style=style)
    final_state = GRAPH.invoke(state)
    if not final_state.get("final_video"):
        raise RuntimeError("최종 동영상 생성에 실패했습니다.")
    return final_state
